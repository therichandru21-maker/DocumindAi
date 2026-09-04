import os
import re
import json
import uuid
import hashlib
import zipfile
from pathlib import Path
from datetime import datetime, timezone
from typing import List, Dict, Any, Optional

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# ---------------------------------------------------------
# OPTIONAL LIBRARIES
# ---------------------------------------------------------

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None

try:
    from google import genai
except Exception:
    genai = None


# ---------------------------------------------------------
# PATHS
# ---------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent

UPLOAD_DIR = BASE_DIR / "uploads"
DATA_DIR = BASE_DIR / "data"
DOCUMENTS_FILE = DATA_DIR / "documents.json"

UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
DATA_DIR.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------
# CONFIG
# ---------------------------------------------------------

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".txt",
    ".zip",
}

MAX_FILE_SIZE = 20 * 1024 * 1024

CHUNK_SIZE = 450
CHUNK_OVERLAP = 70

MAX_SEARCH_RESULTS = 6
MAX_FIELD_RESULTS = 6

MIN_SEARCH_SCORE = 0.08

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
GEMINI_MODEL = os.getenv(
    "GEMINI_MODEL",
    "gemini-2.5-flash"
)


# ---------------------------------------------------------
# APP
# ---------------------------------------------------------

app = FastAPI(
    title="DocuMind AI",
    description="Intelligent Document Knowledge Assistant",
    version="3.0.0",
)


# ---------------------------------------------------------
# CORS
# ---------------------------------------------------------

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5173",
        "http://127.0.0.1:5173",
        "http://localhost:3000",
        "http://127.0.0.1:3000",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ---------------------------------------------------------
# GEMINI
# ---------------------------------------------------------

gemini_client = None

if genai and GEMINI_API_KEY:
    try:
        gemini_client = genai.Client(
            api_key=GEMINI_API_KEY
        )
        print("Gemini AI initialized")
    except Exception as exc:
        print("Gemini initialization failed:", exc)
        gemini_client = None


# ---------------------------------------------------------
# MODELS
# ---------------------------------------------------------

class AskRequest(BaseModel):
    question: str


class SearchRequest(BaseModel):
    query: str


# ---------------------------------------------------------
# FIELD ALIASES
# ---------------------------------------------------------

FIELD_ALIASES = {

    "name": [
        "name",
        "student name",
        "candidate name",
        "full name",
        "employee name",
        "customer name",
    ],

    "student id": [
        "student id",
        "studentid",
        "student no",
        "student number",
        "roll no",
        "roll number",
        "register number",
        "registration number",
        "reg no",
        "reg number",
    ],

    "class": [
        "class",
        "class name",
        "class section",
        "section",
        "class/section",
    ],

    "course": [
        "course",
        "course name",
        "program",
        "programme",
        "degree",
        "department",
    ],

    "term": [
        "term",
        "term name",
        "semester",
        "academic term",
        "academic year",
    ],

    "transaction id": [
        "tran id",
        "transaction id",
        "transaction no",
        "transaction number",
        "txn id",
        "txn number",
    ],

    "pg transaction id": [
        "pg tran id",
        "pg transaction id",
        "pg txn id",
        "pg transaction number",
    ],

    "date": [
        "date",
        "transaction date",
        "trans date",
        "trans. date",
        "payment date",
        "issued date",
        "issue date",
    ],

    "time": [
        "time",
        "transaction time",
        "payment time",
        "issued time",
    ],

    "amount": [
        "amount",
        "amount paid",
        "paid amount",
        "payment amount",
    ],

    "fees": [
        "fees",
        "fee",
        "seminar fees",
        "seminar fee",
        "tuition fees",
        "tuition fee",
    ],

    "total": [
        "total",
        "total amount",
        "grand total",
        "total paid",
        "total fees",
        "amount payable",
        "net amount",
    ],

    "email": [
        "email",
        "email address",
        "mail",
        "mail id",
    ],

    "phone": [
        "phone",
        "phone number",
        "mobile",
        "mobile number",
        "contact number",
        "contact no",
    ],

    "address": [
        "address",
        "home address",
        "permanent address",
        "communication address",
    ],

    "status": [
        "status",
        "payment status",
        "application status",
        "result",
    ],
}


# ---------------------------------------------------------
# QUESTION NOISE
# ---------------------------------------------------------

QUESTION_NOISE = {
    "what",
    "is",
    "the",
    "a",
    "an",
    "of",
    "my",
    "me",
    "tell",
    "give",
    "show",
    "find",
    "search",
    "please",
    "can",
    "you",
    "could",
    "would",
    "get",
    "from",
    "document",
    "file",
    "this",
    "that",
    "about",
    "for",
    "in",
    "on",
}


# ---------------------------------------------------------
# DOCUMENT STORAGE
# ---------------------------------------------------------

def load_documents() -> List[Dict[str, Any]]:

    if not DOCUMENTS_FILE.exists():
        return []

    try:
        with open(
            DOCUMENTS_FILE,
            "r",
            encoding="utf-8"
        ) as f:

            data = json.load(f)

        if isinstance(data, list):
            return data

        return []

    except Exception as exc:

        print(
            "Could not load documents:",
            exc
        )

        return []


def save_documents(
    documents: List[Dict[str, Any]]
) -> None:

    temp_file = DOCUMENTS_FILE.with_suffix(".tmp")

    with open(
        temp_file,
        "w",
        encoding="utf-8"
    ) as f:

        json.dump(
            documents,
            f,
            indent=2,
            ensure_ascii=False
        )

    temp_file.replace(DOCUMENTS_FILE)


# ---------------------------------------------------------
# TEXT HELPERS
# ---------------------------------------------------------

def clean_text(text: str) -> str:

    if not text:
        return ""

    text = text.replace(
        "\x00",
        " "
    )

    text = re.sub(
        r"\s+",
        " ",
        text
    )

    return text.strip()


def normalize_text(text: str) -> str:

    if not text:
        return ""

    text = text.lower()

    text = re.sub(
        r"[^a-z0-9\s]",
        " ",
        text
    )

    text = re.sub(
        r"\s+",
        " ",
        text
    )

    return text.strip()


def normalize_label(text: str) -> str:

    return normalize_text(text)


def tokenize(text: str) -> List[str]:

    normalized = normalize_text(text)

    words = normalized.split()

    return [
        word
        for word in words
        if len(word) > 1
    ]


# ---------------------------------------------------------
# EXTRACT PDF
# ---------------------------------------------------------

def extract_pdf(path: Path) -> str:

    if PdfReader is None:
        raise RuntimeError(
            "pypdf is not installed"
        )

    reader = PdfReader(str(path))

    pages = []

    for page in reader.pages:

        try:

            text = page.extract_text()

            if text:
                pages.append(text)

        except Exception as exc:

            print(
                "PDF page extraction error:",
                exc
            )

    return "\n".join(pages)


# ---------------------------------------------------------
# EXTRACT DOCX
# ---------------------------------------------------------

def extract_docx(path: Path) -> str:

    if Document is None:
        raise RuntimeError(
            "python-docx is not installed"
        )

    document = Document(str(path))

    parts = []

    for paragraph in document.paragraphs:

        if paragraph.text.strip():

            parts.append(
                paragraph.text
            )

    for table in document.tables:

        for row in table.rows:

            values = []

            for cell in row.cells:

                values.append(
                    cell.text.strip()
                )

            parts.append(
                " | ".join(values)
            )

    return "\n".join(parts)


# ---------------------------------------------------------
# EXTRACT PPTX
# ---------------------------------------------------------

def extract_pptx(path: Path) -> str:

    if Presentation is None:
        raise RuntimeError(
            "python-pptx is not installed"
        )

    presentation = Presentation(
        str(path)
    )

    slides = []

    for slide in presentation.slides:

        slide_parts = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:
                    slide_parts.append(text)

        if slide_parts:

            slides.append(
                "\n".join(slide_parts)
            )

    return "\n".join(slides)


# ---------------------------------------------------------
# EXTRACT XLSX
# ---------------------------------------------------------

def extract_xlsx(path: Path) -> str:

    if load_workbook is None:
        raise RuntimeError(
            "openpyxl is not installed"
        )

    workbook = load_workbook(
        filename=str(path),
        read_only=True,
        data_only=True
    )

    parts = []

    for sheet in workbook.worksheets:

        parts.append(
            f"Sheet: {sheet.title}"
        )

        for row in sheet.iter_rows(
            values_only=True
        ):

            values = []

            for value in row:

                if value is not None:

                    values.append(
                        str(value)
                    )

            if values:

                parts.append(
                    " | ".join(values)
                )

    return "\n".join(parts)


# ---------------------------------------------------------
# EXTRACT TXT
# ---------------------------------------------------------

def extract_txt(path: Path) -> str:

    encodings = [
        "utf-8",
        "utf-8-sig",
        "latin-1",
    ]

    for encoding in encodings:

        try:

            return path.read_text(
                encoding=encoding
            )

        except Exception:
            continue

    return ""


# ---------------------------------------------------------
# EXTRACT ZIP
# ---------------------------------------------------------

def extract_zip(path: Path) -> str:

    parts = []

    allowed_text_extensions = {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".py",
        ".js",
        ".jsx",
        ".ts",
        ".tsx",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".css",
        ".html",
        ".sql",
        ".xml",
        ".yaml",
        ".yml",
    }

    try:

        with zipfile.ZipFile(
            path,
            "r"
        ) as archive:

            for info in archive.infolist():

                if info.is_dir():
                    continue

                member_path = Path(
                    info.filename
                )

                extension = (
                    member_path.suffix.lower()
                )

                if extension not in allowed_text_extensions:
                    continue

                if info.file_size > 5 * 1024 * 1024:
                    continue

                try:

                    data = archive.read(
                        info
                    )

                    text = data.decode(
                        "utf-8",
                        errors="ignore"
                    )

                    if text.strip():

                        parts.append(
                            f"File: {info.filename}\n{text}"
                        )

                except Exception as exc:

                    print(
                        "ZIP member error:",
                        exc
                    )

    except zipfile.BadZipFile:

        raise RuntimeError(
            "Invalid ZIP file"
        )

    return "\n".join(parts)


# ---------------------------------------------------------
# DOCUMENT EXTRACTION
# ---------------------------------------------------------

def extract_document(
    path: Path
) -> str:

    extension = (
        path.suffix.lower()
    )

    if extension == ".pdf":

        return extract_pdf(path)

    if extension == ".docx":

        return extract_docx(path)

    if extension == ".pptx":

        return extract_pptx(path)

    if extension == ".xlsx":

        return extract_xlsx(path)

    if extension == ".txt":

        return extract_txt(path)

    if extension == ".zip":

        return extract_zip(path)

    raise ValueError(
        f"Unsupported file type: {extension}"
    )


# ---------------------------------------------------------
# CHUNKING
# ---------------------------------------------------------

def create_chunks(
    text: str,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> List[str]:

    text = clean_text(text)

    words = text.split()

    chunks = []

    start = 0

    while start < len(words):

        end = min(
            start + chunk_size,
            len(words)
        )

        chunk = " ".join(
            words[start:end]
        ).strip()

        if chunk:

            chunks.append(
                chunk
            )

        if end >= len(words):
            break

        start = max(
            end - overlap,
            start + 1
        )

    return chunks


# ---------------------------------------------------------
# FIELD DETECTION
# ---------------------------------------------------------

def detect_requested_fields(
    question: str
) -> List[str]:

    normalized_question = normalize_text(
        question
    )

    found = []

    for canonical, aliases in FIELD_ALIASES.items():

        best_match = None

        for alias in aliases:

            alias_normalized = normalize_text(
                alias
            )

            if not alias_normalized:
                continue

            pattern = (
                rf"\b{re.escape(alias_normalized)}\b"
            )

            if re.search(
                pattern,
                normalized_question
            ):

                if (
                    best_match is None
                    or len(alias_normalized)
                    > len(best_match)
                ):

                    best_match = alias_normalized

        if best_match:

            found.append(
                (
                    canonical,
                    best_match
                )
            )

    found.sort(
        key=lambda item: len(item[1]),
        reverse=True
    )

    return [
        item[0]
        for item in found
    ]


# ---------------------------------------------------------
# FIELD PAIR EXTRACTION
# ---------------------------------------------------------

def extract_field_pairs(
    text: str
) -> List[Dict[str, str]]:

    if not text:
        return []

    aliases = []

    for field_aliases in FIELD_ALIASES.values():

        aliases.extend(
            field_aliases
        )

    aliases = sorted(
        set(
            normalize_text(alias)
            for alias in aliases
            if alias
        ),
        key=len,
        reverse=True
    )

    if not aliases:
        return []

    escaped_aliases = []

    for alias in aliases:

        escaped = re.escape(alias)

        escaped = escaped.replace(
            r"\ ",
            r"\s+"
        )

        escaped_aliases.append(
            escaped
        )

    label_pattern = "|".join(
        escaped_aliases
    )

    pattern = re.compile(
        rf"""
        (?<![A-Za-z0-9])
        (?P<label>{label_pattern})
        [\s.()/_-]*
        [:=]
        \s*
        (?P<value>.*?)
        (?=
            \s+
            (?:
                {label_pattern}
            )
            [\s.()/_-]*
            [:=]
            |
            $
        )
        """,
        re.IGNORECASE | re.VERBOSE
    )

    pairs = []

    # Keep lines where possible
    lines = re.split(
        r"[\r\n]+",
        text
    )

    # Also scan complete flattened text
    # because PDF extraction often removes line breaks.
    scan_texts = list(lines)

    flattened = clean_text(text)

    if flattened:
        scan_texts.append(
            flattened
        )

    for scan_text in scan_texts:

        if not scan_text.strip():
            continue

        for match in pattern.finditer(
            scan_text
        ):

            label = clean_text(
                match.group("label")
            )

            value = clean_text(
                match.group("value")
            )

            if not label or not value:
                continue

            value = re.sub(
                r"^[\s:=-]+",
                "",
                value
            )

            value = re.sub(
                r"[\s|]+$",
                "",
                value
            )

            if value:

                pairs.append(
                    {
                        "label": label,
                        "value": value,
                    }
                )

    # Remove duplicates
    unique = []

    seen = set()

    for pair in pairs:

        key = (
            normalize_label(pair["label"]),
            normalize_text(pair["value"])
        )

        if key in seen:
            continue

        seen.add(key)

        unique.append(pair)

    return unique


# ---------------------------------------------------------
# FIELD MATCHING
# ---------------------------------------------------------

def find_field_values(
    question: str,
    documents: List[Dict[str, Any]]
) -> List[Dict[str, Any]]:

    requested_fields = detect_requested_fields(
        question
    )

    if not requested_fields:
        return []

    results = []

    for document in documents:

        source_text = (
            document.get("raw_text")
            or document.get("text")
            or ""
        )

        pairs = extract_field_pairs(
            source_text
        )

        if not pairs:
            continue

        for field in requested_fields:

            aliases = {
                normalize_label(alias)
                for alias in FIELD_ALIASES[field]
            }

            canonical_normalized = normalize_label(
                field
            )

            matching_pairs = []

            for pair in pairs:

                label_normalized = normalize_label(
                    pair["label"]
                )

                if label_normalized in aliases:

                    matching_pairs.append(
                        pair
                    )

            if not matching_pairs:
                continue

            # Prefer exact canonical label
            matching_pairs.sort(
                key=lambda pair: (
                    0
                    if normalize_label(
                        pair["label"]
                    ) == canonical_normalized
                    else 1,
                    -len(
                        normalize_label(
                            pair["label"]
                        )
                    ),
                )
            )

            best_pair = matching_pairs[0]

            results.append(
                {
                    "id": document["id"],
                    "filename": document["filename"],
                    "extension": document["extension"],
                    "text": best_pair["value"],
                    "field": field,
                    "label": best_pair["label"],
                    "score": 1.0,
                    "chunk_index": -1,
                }
            )

    # Remove duplicates
    unique_results = []

    seen = set()

    for result in results:

        key = (
            result["filename"],
            result["field"],
            normalize_text(
                result["text"]
            )
        )

        if key in seen:
            continue

        seen.add(key)

        unique_results.append(
            result
        )

    return unique_results[
        :MAX_FIELD_RESULTS
    ]


# ---------------------------------------------------------
# SEARCH SCORING
# ---------------------------------------------------------

def calculate_search_score(
    query: str,
    text: str
) -> float:

    query_normalized = normalize_text(
        query
    )

    text_normalized = normalize_text(
        text
    )

    if not query_normalized:
        return 0.0

    if not text_normalized:
        return 0.0

    query_tokens = set(
        tokenize(query)
    )

    text_tokens = set(
        tokenize(text)
    )

    if not query_tokens:
        return 0.0

    overlap = (
        len(query_tokens & text_tokens)
        / len(query_tokens)
    )

    exact_match = (
        1.0
        if query_normalized
        in text_normalized
        else 0.0
    )

    phrase_match = 0.0

    if len(query_tokens) >= 2:

        query_words = query_normalized.split()

        if all(
            word in text_normalized
            for word in query_words
        ):

            phrase_match = 1.0

    score = (
        overlap * 0.55
        + exact_match * 0.30
        + phrase_match * 0.15
    )

    return min(
        score,
        1.0
    )


# ---------------------------------------------------------
# SNIPPET
# ---------------------------------------------------------

def create_search_snippet(
    text: str,
    query: str,
    max_length: int = 450
) -> str:

    text = clean_text(text)

    if len(text) <= max_length:
        return text

    normalized_text = text.lower()
    normalized_query = query.lower()

    position = normalized_text.find(
        normalized_query
    )

    if position == -1:

        query_tokens = tokenize(
            query
        )

        position = -1

        for token in query_tokens:

            position = normalized_text.find(
                token.lower()
            )

            if position != -1:
                break

    if position == -1:

        return text[
            :max_length
        ]

    start = max(
        0,
        position - 150
    )

    end = min(
        len(text),
        start + max_length
    )

    snippet = text[
        start:end
    ]

    if start > 0:
        snippet = "..." + snippet

    if end < len(text):
        snippet += "..."

    return snippet


# ---------------------------------------------------------
# NORMAL SEARCH
# ---------------------------------------------------------

def search_documents(
    query: str,
    documents: List[Dict[str, Any]]
) -> List[Dict[str, Any]]:

    query = query.strip()

    if not query:
        return []

    results = []

    for document in documents:

        chunks = document.get(
            "chunks_data",
            []
        )

        # Support older document structure
        if not chunks:

            text = document.get(
                "text",
                ""
            )

            chunks = [
                {
                    "text": text,
                    "chunk_index": 0
                }
            ]

        for chunk in chunks:

            chunk_text = chunk.get(
                "text",
                ""
            )

            if not chunk_text:
                continue

            score = calculate_search_score(
                query,
                chunk_text
            )

            if score < MIN_SEARCH_SCORE:
                continue

            results.append(
                {
                    "id": document["id"],
                    "filename": document["filename"],
                    "extension": document["extension"],
                    "text": create_search_snippet(
                        chunk_text,
                        query
                    ),
                    "score": round(
                        score,
                        4
                    ),
                    "chunk_index": chunk.get(
                        "chunk_index",
                        0
                    ),
                }
            )

    results.sort(
        key=lambda item: item["score"],
        reverse=True
    )

    # Remove duplicates
    final_results = []

    seen = set()

    file_counts = {}

    for result in results:

        normalized_snippet = normalize_text(
            result["text"]
        )[:300]

        key = (
            result["filename"],
            normalized_snippet
        )

        if key in seen:
            continue

        count = file_counts.get(
            result["filename"],
            0
        )

        if count >= 3:
            continue

        seen.add(key)

        file_counts[
            result["filename"]
        ] = count + 1

        final_results.append(
            result
        )

        if len(final_results) >= MAX_SEARCH_RESULTS:
            break

    return final_results


# ---------------------------------------------------------
# GEMINI ANSWER
# ---------------------------------------------------------

def generate_ai_answer(
    question: str,
    search_results: List[Dict[str, Any]]
) -> Optional[str]:

    if not gemini_client:
        return None

    if not search_results:
        return None

    context_parts = []

    for index, result in enumerate(
        search_results[:5],
        start=1
    ):

        context_parts.append(
            f"""
SOURCE {index}
File: {result["filename"]}
Relevant content:
{result["text"]}
"""
        )

    context = "\n".join(
        context_parts
    )

    prompt = f"""
You are DocuMind AI, a document question-answering assistant.

USER QUESTION:
{question}

RELEVANT DOCUMENT INFORMATION:
{context}

STRICT RULES:

1. Answer ONLY the user's question.
2. Use ONLY the information provided in the relevant document information.
3. Do NOT use outside knowledge.
4. Do NOT reproduce the entire document.
5. Do NOT dump the retrieved context.
6. Do NOT mention irrelevant fields.
7. If the question asks for a specific value, return only that value.
8. If the question asks for an explanation, give a short and relevant explanation.
9. If the question asks for a summary, provide a concise summary.
10. If the answer is not available in the provided information, say:
   "I couldn't find that information in the uploaded documents."
11. Keep the answer concise and directly relevant.
12. Never invent or guess information.

ANSWER:
"""

    try:

        response = gemini_client.models.generate_content(
            model=GEMINI_MODEL,
            contents=prompt,
        )

        answer = getattr(
            response,
            "text",
            None
        )

        if answer:

            answer = answer.strip()

            if answer:
                return answer

    except Exception as exc:

        print(
            "Gemini error:",
            exc
        )

    return None


# ---------------------------------------------------------
# FALLBACK ANSWER
# ---------------------------------------------------------

def fallback_answer(
    question: str,
    results: List[Dict[str, Any]]
) -> str:

    if not results:

        return (
            "I couldn't find that information "
            "in the uploaded documents."
        )

    best = results[0]

    return best["text"]


# ---------------------------------------------------------
# ROOT
# ---------------------------------------------------------

@app.get("/")
def root():

    return {
        "app": "DocuMind AI",
        "version": "3.0.0",
        "status": "running",
    }


# ---------------------------------------------------------
# STATUS
# ---------------------------------------------------------

@app.get("/status")
def status():

    documents = load_documents()

    return {
        "ready": True,
        "indexed_files": len(documents),
        "supported_extensions": sorted(
            SUPPORTED_EXTENSIONS
        ),
        "gemini_enabled": (
            gemini_client is not None
        ),
    }


# ---------------------------------------------------------
# DOCUMENT LIST
# ---------------------------------------------------------

@app.get("/documents")
def get_documents():

    documents = load_documents()

    output = []

    for document in documents:

        output.append(
            {
                "id": document["id"],
                "filename": document["filename"],
                "extension": document["extension"],
                "size": document.get(
                    "size",
                    0
                ),
                "created_at": document.get(
                    "created_at"
                ),
                "chunks": len(
                    document.get(
                        "chunks_data",
                        []
                    )
                ),
            }
        )

    return {
        "documents": output
    }


# ---------------------------------------------------------
# UPLOAD
# ---------------------------------------------------------

@app.post("/upload")
@app.post("/documents/upload")
async def upload_document(
    file: UploadFile = File(...)
):

    if not file.filename:

        raise HTTPException(
            status_code=400,
            detail="Filename is required."
        )

    original_filename = Path(
        file.filename
    ).name

    extension = Path(
        original_filename
    ).suffix.lower()

    if extension not in SUPPORTED_EXTENSIONS:

        raise HTTPException(
            status_code=400,
            detail=(
                "Unsupported file type. "
                "Supported: "
                + ", ".join(
                    sorted(
                        SUPPORTED_EXTENSIONS
                    )
                )
            ),
        )

    content = await file.read()

    if not content:

        raise HTTPException(
            status_code=400,
            detail="Uploaded file is empty."
        )

    if len(content) > MAX_FILE_SIZE:

        raise HTTPException(
            status_code=413,
            detail="File size exceeds 20 MB."
        )

    file_hash = hashlib.sha256(
        content
    ).hexdigest()

    documents = load_documents()

    # Prevent exact duplicate upload
    for existing in documents:

        if existing.get(
            "sha256"
        ) == file_hash:

            return {
                "success": True,
                "already_indexed": True,
                "document": {
                    "id": existing["id"],
                    "filename": existing["filename"],
                    "extension": existing["extension"],
                    "chunks": len(
                        existing.get(
                            "chunks_data",
                            []
                        )
                    ),
                },
            }

    document_id = uuid.uuid4().hex

    stored_filename = (
        f"{document_id}{extension}"
    )

    stored_path = (
        UPLOAD_DIR / stored_filename
    )

    try:

        with open(
            stored_path,
            "wb"
        ) as f:

            f.write(content)

        extracted_text = extract_document(
            stored_path
        )

        if not extracted_text.strip():

            stored_path.unlink(
                missing_ok=True
            )

            raise HTTPException(
                status_code=422,
                detail=(
                    "No readable text "
                    "could be extracted "
                    "from this document."
                ),
            )

        chunks = create_chunks(
            extracted_text
        )

        chunks_data = []

        for index, chunk in enumerate(
            chunks
        ):

            chunks_data.append(
                {
                    "chunk_index": index,
                    "text": chunk,
                }
            )

        document = {
            "id": document_id,
            "filename": original_filename,
            "extension": extension,
            "stored_filename": stored_filename,
            "size": len(content),
            "sha256": file_hash,
            "created_at": datetime.now(
                timezone.utc
            ).isoformat(),
            "text": clean_text(
                extracted_text
            ),
            "raw_text": extracted_text,
            "chunks_data": chunks_data,
        }

        documents.append(
            document
        )

        save_documents(
            documents
        )

        return {
            "success": True,
            "already_indexed": False,
            "document": {
                "id": document_id,
                "filename": original_filename,
                "extension": extension,
                "chunks": len(
                    chunks_data
                ),
            },
        }

    except HTTPException:
        raise

    except Exception as exc:

        stored_path.unlink(
            missing_ok=True
        )

        print(
            "Upload error:",
            exc
        )

        raise HTTPException(
            status_code=500,
            detail=(
                f"Failed to process "
                f"document: {str(exc)}"
            ),
        )


# ---------------------------------------------------------
# DELETE DOCUMENT
# ---------------------------------------------------------

@app.delete("/documents/{document_id}")
def delete_document(
    document_id: str
):

    documents = load_documents()

    target = None

    for document in documents:

        if document["id"] == document_id:

            target = document
            break

    if target is None:

        raise HTTPException(
            status_code=404,
            detail="Document not found."
        )

    documents = [
        document
        for document in documents
        if document["id"] != document_id
    ]

    save_documents(
        documents
    )

    stored_filename = target.get(
        "stored_filename"
    )

    if stored_filename:

        stored_path = (
            UPLOAD_DIR / stored_filename
        )

        stored_path.unlink(
            missing_ok=True
        )

    return {
        "success": True,
        "message": (
            "Document deleted successfully."
        ),
    }


# ---------------------------------------------------------
# SEARCH
# ---------------------------------------------------------

@app.post("/search")
@app.post("/documents/search")
def search(request: SearchRequest):

    query = request.query.strip()

    if not query:

        return {
            "query": query,
            "results": [],
        }

    documents = load_documents()

    # IMPORTANT:
    # Exact field extraction happens FIRST.
    field_results = find_field_values(
        query,
        documents
    )

    if field_results:

        return {
            "query": query,
            "results": field_results,
            "mode": "exact_field",
        }

    results = search_documents(
        query,
        documents
    )

    return {
        "query": query,
        "results": results,
        "mode": "semantic_search",
    }


# ---------------------------------------------------------
# ASK DOCUMIND AI
# ---------------------------------------------------------

@app.post("/ask")
@app.post("/documents/ask")
def ask_documind(
    request: AskRequest
):

    question = request.question.strip()

    if not question:

        raise HTTPException(
            status_code=400,
            detail="Question cannot be empty."
        )

    documents = load_documents()

    # =====================================================
    # STEP 1
    # EXACT FIELD EXTRACTION
    # =====================================================

    field_results = find_field_values(
        question,
        documents
    )

    if field_results:

        unique_fields = []

        for result in field_results:

            if result["field"] not in unique_fields:

                unique_fields.append(
                    result["field"]
                )

        # -------------------------------------------------
        # SINGLE FIELD + SINGLE RESULT
        # RETURN ONLY THE VALUE
        # -------------------------------------------------

        if (
            len(unique_fields) == 1
            and len(field_results) == 1
        ):

            answer = field_results[0][
                "text"
            ]

        # -------------------------------------------------
        # MULTIPLE FIELDS
        # -------------------------------------------------

        else:

            answer_parts = []

            used = set()

            for result in field_results:

                key = (
                    result["field"],
                    result["text"]
                )

                if key in used:
                    continue

                used.add(key)

                readable_field = (
                    result["field"]
                    .title()
                )

                answer_parts.append(
                    f"{readable_field}: "
                    f"{result['text']}"
                )

            answer = "\n".join(
                answer_parts
            )

        sources = list(
            dict.fromkeys(
                result["filename"]
                for result in field_results
            )
        )

        return {
            "question": question,
            "answer": answer,
            "sources": sources,
            "results": field_results,
            "mode": "exact_field",
        }

    # =====================================================
    # STEP 2
    # NORMAL DOCUMENT RETRIEVAL
    # =====================================================

    search_results = search_documents(
        question,
        documents
    )

    if not search_results:

        return {
            "question": question,
            "answer": (
                "I couldn't find that information "
                "in the uploaded documents."
            ),
            "sources": [],
            "results": [],
            "mode": "no_match",
        }

    # =====================================================
    # STEP 3
    # GEMINI ANSWERING
    # =====================================================

    ai_answer = generate_ai_answer(
        question,
        search_results
    )

    if ai_answer:

        answer = ai_answer
        mode = "ai_rag"

    else:

        answer = fallback_answer(
            question,
            search_results
        )

        mode = "retrieval_fallback"

    sources = list(
        dict.fromkeys(
            result["filename"]
            for result in search_results
        )
    )

    return {
        "question": question,
        "answer": answer,
        "sources": sources,
        "results": search_results,
        "mode": mode,
    }


# ---------------------------------------------------------
# RUN
# ---------------------------------------------------------

if __name__ == "__main__":

    import uvicorn

    uvicorn.run(
        "main:app",
        host="127.0.0.1",
        port=8000,
        reload=True,
    )